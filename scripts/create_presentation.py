"""Generate course project presentation (PPTX) for Environmental Health Intelligence Platform."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Colors ───────────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT_BLUE = RGBColor(0x16, 0x21, 0x3E)   # mid navy
TEAL        = RGBColor(0x0F, 0x3D, 0x4A)   # dark teal
HIGHLIGHT   = RGBColor(0x00, 0xB4, 0xD8)   # cyan/teal highlight
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # warning orange
GREEN       = RGBColor(0x2D, 0xC6, 0x71)   # positive green
RED_LIGHT   = RGBColor(0xFF, 0x4D, 0x4D)   # danger red
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xDD, 0xEE)
MID_GRAY    = RGBColor(0x88, 0x99, 0xAA)
YELLOW      = RGBColor(0xFF, 0xD6, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()   # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  font_size=16, color=WHITE, bold=False,
                  line_spacing=1.15, align=PP_ALIGN.LEFT):
    """lines: list of (text, bold_override, color_override, size_override)"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, b, c, s = item, bold, color, font_size
        else:
            text = item[0]
            b    = item[1] if len(item) > 1 else bold
            c    = item[2] if len(item) > 2 else color
            s    = item[3] if len(item) > 3 else font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
    return txBox


def add_divider(slide, t, color=HIGHLIGHT, w=13.0, l=0.165):
    rect = add_rect(slide, l, t, w, 0.04, fill_color=color)
    return rect


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)

    # left accent bar
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=HIGHLIGHT)

    # top decoration rectangle
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    # main title
    add_text(s, "Environmental Health Intelligence Platform",
             0.5, 1.2, 12.3, 2.0,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # subtitle
    add_text(s, "How Air Pollution Kills — and What the Data Says We Should Do About It",
             0.5, 3.0, 12.0, 1.0,
             font_size=22, bold=False, color=HIGHLIGHT, align=PP_ALIGN.LEFT)

    add_divider(s, 4.1)

    # metadata row
    add_multiline(s, [
        ("Data Sources: WHO GHO · IHME GBD · OpenAQ v3 · Open-Meteo ERA5 · World Bank",
         False, LIGHT_GRAY, 15),
        ("Countries covered: 183+   |   Years: 2010 – 2024   |   Tool: Python · Plotly · Tableau",
         False, MID_GRAY, 13),
    ], 0.5, 4.25, 12.0, 0.9)

    # bottom tag
    add_text(s, "Oleg Sher  |  Data & Business Analysis Course  |  2026",
             0.5, 6.8, 12.0, 0.5,
             font_size=13, color=MID_GRAY, align=PP_ALIGN.LEFT)
    return s


def slide_about_me(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=HIGHLIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Who Am I?", 0.5, 0.3, 12.0, 0.7,
             font_size=32, bold=True, color=WHITE)
    add_divider(s, 1.15)

    # icon circle placeholder
    add_rect(s, 0.5, 1.4, 3.5, 4.5, fill_color=ACCENT_BLUE)
    add_text(s, "👤", 1.6, 2.5, 1.5, 1.5, font_size=60, align=PP_ALIGN.CENTER)

    bullets = [
        ("Oleg Sher", True, WHITE, 24),
        ("", False, WHITE, 8),
        ("Data & Business Analyst", False, HIGHLIGHT, 18),
        ("", False, WHITE, 8),
        ("• Background in data engineering & analytics", False, LIGHT_GRAY, 16),
        ("• Passionate about turning raw data into actionable insights", False, LIGHT_GRAY, 16),
        ("• This project merges environmental science with business intelligence", False, LIGHT_GRAY, 16),
        ("", False, WHITE, 8),
        ("• Tools: Python, Pandas, Plotly, Tableau, Kaggle", False, MID_GRAY, 14),
        ("• Datasets published on Kaggle for public research use", False, MID_GRAY, 14),
    ]
    add_multiline(s, bullets, 4.3, 1.4, 8.7, 5.0, font_size=16, color=WHITE)
    return s


def slide_research_question(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=ORANGE)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Research Question & Dataset", 0.5, 0.3, 12.0, 0.7,
             font_size=32, bold=True, color=WHITE)
    add_divider(s, 1.15, color=ORANGE)

    # Big question box
    add_rect(s, 0.5, 1.35, 12.3, 1.4, fill_color=TEAL)
    add_text(s,
             "How does air pollution (PM2.5) impact human mortality across countries — "
             "and what levers can policymakers pull to reduce the burden?",
             0.65, 1.45, 12.0, 1.2,
             font_size=19, bold=True, color=HIGHLIGHT, align=PP_ALIGN.LEFT)

    # 3-column metadata
    cols = [
        ("Topic", "Air quality (PM2.5 / PM10)\n& health burden", HIGHLIGHT),
        ("Time period", "2010 – 2024\n(WHO / IHME up to 2023)", ORANGE),
        ("Geography", "183+ countries\nglobal coverage", GREEN),
    ]
    for i, (label, value, c) in enumerate(cols):
        xl = 0.5 + i * 4.25
        add_rect(s, xl, 2.95, 4.0, 1.6, fill_color=ACCENT_BLUE)
        add_text(s, label, xl + 0.15, 3.05, 3.7, 0.4,
                 font_size=13, color=c, bold=True)
        add_text(s, value, xl + 0.15, 3.45, 3.7, 0.9,
                 font_size=15, color=WHITE)

    # Data sources row
    add_text(s, "Data Sources", 0.5, 4.75, 2.5, 0.4,
             font_size=14, bold=True, color=MID_GRAY)
    sources = [
        ("WHO GHO", "PM2.5 concentration &\nmortality indicators"),
        ("IHME GBD", "Death rates attributable\nto air pollution"),
        ("OpenAQ v3", "Real-time station readings\n8 key countries"),
        ("World Bank", "GDP, population,\nlife expectancy"),
    ]
    for i, (name, desc) in enumerate(sources):
        xl = 0.5 + i * 3.2
        add_rect(s, xl, 5.15, 3.0, 1.6, fill_color=RGBColor(0x0D, 0x2B, 0x3E))
        add_text(s, name, xl + 0.12, 5.22, 2.8, 0.4,
                 font_size=13, bold=True, color=HIGHLIGHT)
        add_text(s, desc, xl + 0.12, 5.62, 2.8, 0.8,
                 font_size=12, color=LIGHT_GRAY)
    return s


def slide_q1_dirtiest(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=RED_LIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Q1 — Where Is the Air Dirtiest?", 0.5, 0.25, 12.0, 0.55,
             font_size=28, bold=True, color=WHITE)
    add_text(s, "Which countries exceed the WHO PM2.5 annual limit the most — and by how much?",
             0.5, 0.85, 12.5, 0.45, font_size=15, color=LIGHT_GRAY)
    add_divider(s, 1.4, color=RED_LIGHT)

    # Chart placeholder
    add_rect(s, 0.5, 1.55, 7.8, 5.3, fill_color=ACCENT_BLUE)
    add_text(s, "[ Plotly Bar Chart ]\nTop 30 Countries by Annual Mean PM2.5\n(WHO GHO, 2023)\n\nRed ≥45 µg/m³  Orange ≥30  Green <30\nBlue dashed line = WHO limit 15 µg/m³",
             0.65, 2.8, 7.5, 2.5, font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Key stats
    stats = [
        ("🔴  BGD / IND / PAK", "3–4× WHO limit", RED_LIGHT),
        ("70%", "of 183 countries\nexceed 15 µg/m³", ORANGE),
        ("9×", "worst-country ratio\nvs WHO guideline", YELLOW),
    ]
    for i, (num, label, c) in enumerate(stats):
        yt = 1.7 + i * 1.65
        add_rect(s, 8.55, yt, 4.5, 1.4, fill_color=ACCENT_BLUE)
        add_text(s, num,   8.7, yt + 0.05, 4.2, 0.6, font_size=20, bold=True, color=c)
        add_text(s, label, 8.7, yt + 0.65, 4.2, 0.7, font_size=13, color=LIGHT_GRAY)

    add_rect(s, 0.5, 6.95, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Insight: Classify countries by 'multiples of WHO limit' (not pass/fail) for tiered risk-based policy response.",
             0.65, 6.98, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_q2_correlation(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=ORANGE)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Q2 — Does More Pollution Mean More Deaths?",
             0.5, 0.25, 12.0, 0.55, font_size=28, bold=True, color=WHITE)
    add_text(s, "Is there a measurable country-level link between PM2.5 and air-pollution mortality?",
             0.5, 0.85, 12.5, 0.45, font_size=15, color=LIGHT_GRAY)
    add_divider(s, 1.4, color=ORANGE)

    add_rect(s, 0.5, 1.55, 7.8, 5.3, fill_color=ACCENT_BLUE)
    add_text(s, "[ Plotly Scatter Chart ]\nPM2.5 vs Mortality Rate — 183 Countries (2019)\n\nColor = Region  Trendline = OLS regression\nDashed = WHO 15 µg/m³ limit",
             0.65, 2.9, 7.5, 2.5, font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    findings = [
        ("r ≈ 0.6", "Positive correlation —\nbut NOT deterministic", HIGHLIGHT),
        ("Sub-Saharan\nAfrica", "High mortality despite\nmoderate outdoor PM2.5\n→ indoor cooking fires", ORANGE),
        ("Gulf states", "Extreme PM2.5, lower\nmortality → income &\nhealthcare buffer", GREEN),
    ]
    for i, (num, label, c) in enumerate(findings):
        yt = 1.7 + i * 1.65
        add_rect(s, 8.55, yt, 4.5, 1.4, fill_color=ACCENT_BLUE)
        add_text(s, num,   8.7, yt + 0.05, 4.2, 0.55, font_size=17, bold=True, color=c)
        add_text(s, label, 8.7, yt + 0.6,  4.2, 0.75, font_size=12, color=LIGHT_GRAY)

    add_rect(s, 0.5, 6.95, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Insight: ROI of pollution reduction differs sharply by region — outdoor monitors miss the indoor biomass burning burden.",
             0.65, 6.98, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_q3_per_capita(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=YELLOW)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Q3 — Who Bears the Highest Per-Capita Death Toll?",
             0.5, 0.25, 12.0, 0.55, font_size=27, bold=True, color=WHITE)
    add_text(s, "Which countries have the most air-pollution deaths per 100,000 people? (IHME GBD, 2015)",
             0.5, 0.85, 12.5, 0.45, font_size=15, color=LIGHT_GRAY)
    add_divider(s, 1.4, color=YELLOW)

    add_rect(s, 0.5, 1.55, 7.8, 5.3, fill_color=ACCENT_BLUE)
    add_text(s, "[ Plotly Horizontal Bar Chart ]\nTop 20 Countries: Air Pollution Death Rate/100k\n\nColor scale = Reds (darker = higher rate)\nIHME GBD / Lelieveld et al. 2019",
             0.65, 2.9, 7.5, 2.5, font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    findings = [
        ("Eastern\nEurope leads", "Bulgaria, Hungary,\nPoland — NOT Middle East", YELLOW),
        ("Why?", "Coal + diesel combustion\nmore toxic than desert dust\nper µg/m³", ORANGE),
        ("Key lever", "EU coal phase-out &\nbuilding retrofit →\nhighest lives saved", GREEN),
    ]
    for i, (num, label, c) in enumerate(findings):
        yt = 1.7 + i * 1.65
        add_rect(s, 8.55, yt, 4.5, 1.4, fill_color=ACCENT_BLUE)
        add_text(s, num,   8.7, yt + 0.05, 4.2, 0.55, font_size=17, bold=True, color=c)
        add_text(s, label, 8.7, yt + 0.6,  4.2, 0.75, font_size=12, color=LIGHT_GRAY)

    add_rect(s, 0.5, 6.95, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Insight: The countries with highest mortality rates are NOT the most polluted — combustion source matters more than raw concentration.",
             0.65, 6.98, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_q4_trends(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=GREEN)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Q4 — Are We Winning or Losing? PM2.5 Trends 2010–2023",
             0.5, 0.25, 12.0, 0.55, font_size=26, bold=True, color=WHITE)
    add_text(s, "In the worst-polluted countries — is the situation improving or deteriorating?",
             0.5, 0.85, 12.5, 0.45, font_size=15, color=LIGHT_GRAY)
    add_divider(s, 1.4, color=GREEN)

    add_rect(s, 0.5, 1.55, 12.3, 3.4, fill_color=ACCENT_BLUE)
    add_text(s, "[ Plotly Dual Panel Chart ]\nLeft: Absolute PM2.5 (µg/m³)  Right: % Change vs 2010\nTop 10 most-polluted countries, 2010–2023\nBlue dashed = WHO 15 µg/m³ limit",
             0.65, 2.5, 12.0, 2.0, font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    cagr_data = [
        ("🔴 BGD", "+2.1%/yr", "Worsening", RED_LIGHT),
        ("🟡 IND", "+0.3%/yr", "Flat", YELLOW),
        ("🟡 PAK", "+0.1%/yr", "Flat", YELLOW),
        ("🟢 CHN", "-1.8%/yr", "Improving", GREEN),
        ("🔴 NPL", "+1.4%/yr", "Worsening", RED_LIGHT),
    ]
    for i, (country, cagr, status, c) in enumerate(cagr_data):
        col = i % 3
        row = i // 3
        xl = 0.5 + col * 4.1
        yt = 5.2 + row * 0.85
        add_rect(s, xl, yt, 3.9, 0.75, fill_color=ACCENT_BLUE)
        add_text(s, f"{country}  {cagr}  — {status}",
                 xl + 0.12, yt + 0.1, 3.7, 0.55, font_size=13, bold=False, color=c)

    add_rect(s, 0.5, 6.95, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Insight: At current CAGR, no top-10 country reaches WHO guideline before 2100. Binding 5-year targets with sanctions are the only proven mechanism.",
             0.65, 6.98, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_q5_exceedance(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=RED_LIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Q5 — How Effective is the WHO Guideline?",
             0.5, 0.25, 12.0, 0.55, font_size=28, bold=True, color=WHITE)
    add_text(s, "In 8 directly-monitored countries: how many days exceeded the WHO PM2.5 daily limit?",
             0.5, 0.85, 12.5, 0.45, font_size=15, color=LIGHT_GRAY)
    add_divider(s, 1.4, color=RED_LIGHT)

    add_rect(s, 0.5, 1.55, 7.8, 5.3, fill_color=ACCENT_BLUE)
    add_text(s, "[ Plotly Dual Chart ]\nLeft: % of readings exceeding WHO limit\nRight: Mean vs Peak PM2.5 scatter\n8 countries — Dec 2024 to May 2026",
             0.65, 2.9, 7.5, 2.5, font_size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    exceed_data = [
        ("Bangladesh", "~100%", RED_LIGHT),
        ("India",      " ~98%", RED_LIGHT),
        ("Pakistan",   " ~97%", RED_LIGHT),
        ("Egypt",      " ~95%", ORANGE),
        ("China",      " ~80%", ORANGE),
        ("Indonesia",  " ~55%", YELLOW),
        ("Israel",     " ~30%", YELLOW),
        ("USA",        " ~17%", GREEN),
    ]
    add_text(s, "% Days Over WHO Limit", 8.55, 1.55, 4.5, 0.4,
             font_size=13, bold=True, color=MID_GRAY)
    for i, (country, pct, c) in enumerate(exceed_data):
        yt = 2.05 + i * 0.6
        add_rect(s, 8.55, yt, 4.5, 0.52, fill_color=ACCENT_BLUE)
        add_text(s, f"{country:<14} {pct}",
                 8.7, yt + 0.05, 4.3, 0.42,
                 font_size=13, color=c, bold=(c == RED_LIGHT))

    add_rect(s, 0.5, 6.95, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Insight: Pakistan peak = 575 µg/m³ (38× WHO). USA (17%) proves 15 µg/m³ is achievable — after decades of Clean Air Act enforcement.",
             0.65, 6.98, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_dashboard(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=HIGHLIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Interactive Dashboard", 0.5, 0.3, 12.0, 0.65,
             font_size=32, bold=True, color=WHITE)
    add_divider(s, 1.1)

    # Dashboard preview box
    add_rect(s, 0.5, 1.25, 8.5, 5.0, fill_color=ACCENT_BLUE)
    add_text(s,
             "[ Tableau Dashboard Preview ]\n\n"
             "• World choropleth: PM2.5 by country\n"
             "• Top 20 most-polluted countries bar chart\n"
             "• PM2.5 vs mortality scatter (interactive filters)\n"
             "• WHO exceedance % by country\n"
             "• Time-series trend: 2010–2023",
             0.65, 2.3, 8.2, 3.5, font_size=14, color=MID_GRAY)

    # URL box
    add_rect(s, 9.2, 1.25, 3.9, 1.2, fill_color=TEAL)
    add_text(s, "Tableau Public", 9.35, 1.3, 3.7, 0.45,
             font_size=14, bold=True, color=HIGHLIGHT)
    add_text(s, "public.tableau.com\n[see project Tableau dashboard]",
             9.35, 1.75, 3.7, 0.65, font_size=12, color=LIGHT_GRAY)

    # Kaggle link
    add_rect(s, 9.2, 2.65, 3.9, 1.2, fill_color=TEAL)
    add_text(s, "Kaggle Dataset", 9.35, 2.7, 3.7, 0.45,
             font_size=14, bold=True, color=HIGHLIGHT)
    add_text(s, "kaggle.com/datasets\nglobal-environmental-intelligence",
             9.35, 3.15, 3.7, 0.65, font_size=12, color=LIGHT_GRAY)

    # Kaggle Notebook link
    add_rect(s, 9.2, 4.05, 3.9, 1.2, fill_color=TEAL)
    add_text(s, "Kaggle Notebook", 9.35, 4.1, 3.7, 0.45,
             font_size=14, bold=True, color=HIGHLIGHT)
    add_text(s, "Air Quality Business\nAnalysis — 5 Questions",
             9.35, 4.55, 3.7, 0.65, font_size=12, color=LIGHT_GRAY)

    add_rect(s, 0.5, 6.4, 12.3, 0.45, fill_color=TEAL)
    add_text(s,
             "Dashboard enables policymakers to filter by region, year, and indicator — drill from global overview to country-level specifics.",
             0.65, 6.43, 12.0, 0.4, font_size=12, color=HIGHLIGHT)
    return s


def slide_conclusions(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=GREEN)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Conclusions", 0.5, 0.3, 12.0, 0.65,
             font_size=32, bold=True, color=WHITE)
    add_divider(s, 1.1, color=GREEN)

    conclusions = [
        ("1", "Air pollution is universal, not an outlier",
         "70% of all tracked countries breach the WHO PM2.5 limit. Risk must be classified by multiples, not binary pass/fail.", RED_LIGHT),
        ("2", "Correlation is real — but source matters",
         "r ≈ 0.6 between PM2.5 and mortality. Indoor biomass burning (Sub-Saharan Africa) and income-buffering (Gulf) create regional divergence.", ORANGE),
        ("3", "Eastern Europe's combustion crisis",
         "Highest per-capita death rates are NOT in the most polluted countries — coal and diesel combustion generate finer, deadlier particles.", YELLOW),
        ("4", "No country is on track to meet WHO targets",
         "Most top-10 polluted countries show <1% annual improvement. At current CAGR, the WHO guideline won't be met before 2100.", RED_LIGHT),
        ("5", "Monitoring without enforcement is theatre",
         "Bangladesh/India/Pakistan exceed the limit on 95–100% of days. The USA (17%) proves compliance is achievable with sustained legal enforcement.", HIGHLIGHT),
    ]
    for i, (num, title, desc, c) in enumerate(conclusions):
        yt = 1.3 + i * 1.1
        add_rect(s, 0.5, yt, 0.65, 0.85, fill_color=c)
        add_text(s, num, 0.5, yt, 0.65, 0.85,
                 font_size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(s, 1.25, yt, 11.5, 0.85, fill_color=ACCENT_BLUE)
        add_text(s, title, 1.4, yt + 0.03, 4.5, 0.4,
                 font_size=14, bold=True, color=c)
        add_text(s, desc, 1.4, yt + 0.43, 11.2, 0.4,
                 font_size=12, color=LIGHT_GRAY)
    return s


def slide_recommendations(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=HIGHLIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Recommendations", 0.5, 0.3, 12.0, 0.65,
             font_size=32, bold=True, color=WHITE)
    add_divider(s, 1.1)

    recs = [
        ("Policy", [
            "Replace pass/fail WHO monitoring with tiered risk scoring (1×/2×/5× WHO)",
            "Mandate binding 5-year reduction targets backed by financial sanctions",
            "Prioritize combustion source controls over general air quality programs",
        ], ORANGE),
        ("Public Health", [
            "Separate indoor and outdoor pollution tracks in national burden estimates",
            "Deploy PM2.5 alert systems in South Asia + Sub-Saharan Africa for episodic peaks",
            "Use IHME death-rate data (not PM2.5 alone) to allocate healthcare resources",
        ], RED_LIGHT),
        ("Data & Technology", [
            "Expand real-time monitoring networks in data-sparse regions (Africa, Central Asia)",
            "Open-source country-level pollution dashboards for citizen accountability",
            "Integrate weather/ERA5 context to separate pollution from natural dust events",
        ], HIGHLIGHT),
    ]
    for i, (category, bullets, c) in enumerate(recs):
        xl = 0.4 + i * 4.3
        add_rect(s, xl, 1.35, 4.1, 0.5, fill_color=c)
        add_text(s, category, xl + 0.1, 1.38, 3.9, 0.44,
                 font_size=16, bold=True, color=DARK_BG)
        add_rect(s, xl, 1.9, 4.1, 4.6, fill_color=ACCENT_BLUE)
        for j, b in enumerate(bullets):
            add_text(s, f"• {b}", xl + 0.12, 2.0 + j * 1.35, 3.88, 1.25,
                     font_size=13, color=LIGHT_GRAY)
    return s


def slide_ab_testing(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=YELLOW)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "A/B Testing — Do Regulations Actually Work?",
             0.5, 0.25, 12.0, 0.6, font_size=28, bold=True, color=WHITE)
    add_divider(s, 1.0, color=YELLOW)

    # Hypothesis
    add_rect(s, 0.5, 1.15, 12.3, 0.75, fill_color=TEAL)
    add_text(s, "H₀: OECD countries (strict environmental regulations) have the same PM2.5 levels as non-OECD countries",
             0.65, 1.2, 12.0, 0.35, font_size=13, color=LIGHT_GRAY)
    add_text(s, "H₁: OECD countries have significantly lower PM2.5 — regulations have measurable impact",
             0.65, 1.55, 12.0, 0.3, font_size=13, bold=True, color=HIGHLIGHT)

    # Group A
    add_rect(s, 0.5, 2.1, 5.9, 4.0, fill_color=ACCENT_BLUE)
    add_text(s, "Group A — OECD Countries (38)", 0.65, 2.15, 5.6, 0.45,
             font_size=16, bold=True, color=GREEN)
    add_multiline(s, [
        ("Mean PM2.5: ~14.2 µg/m³", False, LIGHT_GRAY, 14),
        ("", False, LIGHT_GRAY, 6),
        ("• Strict air quality legislation", False, LIGHT_GRAY, 13),
        ("• EU ETS, Clean Air Act (US)", False, LIGHT_GRAY, 13),
        ("• Regular enforcement & penalties", False, LIGHT_GRAY, 13),
        ("• High monitoring coverage", False, LIGHT_GRAY, 13),
        ("", False, LIGHT_GRAY, 6),
        ("Exceedance rate: ~25%", False, GREEN, 14),
    ], 0.65, 2.65, 5.6, 3.2, font_size=13, color=LIGHT_GRAY)

    # VS divider
    add_text(s, "VS", 6.5, 3.4, 0.8, 0.7,
             font_size=26, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Group B
    add_rect(s, 7.4, 2.1, 5.9, 4.0, fill_color=ACCENT_BLUE)
    add_text(s, "Group B — Non-OECD Countries (145+)", 7.55, 2.15, 5.6, 0.45,
             font_size=16, bold=True, color=RED_LIGHT)
    add_multiline(s, [
        ("Mean PM2.5: ~31.6 µg/m³", False, LIGHT_GRAY, 14),
        ("", False, LIGHT_GRAY, 6),
        ("• Limited or unenforced regulations", False, LIGHT_GRAY, 13),
        ("• Rapid industrialization", False, LIGHT_GRAY, 13),
        ("• Sparse monitoring coverage", False, LIGHT_GRAY, 13),
        ("• Agriculture burning, coal use", False, LIGHT_GRAY, 13),
        ("", False, LIGHT_GRAY, 6),
        ("Exceedance rate: ~85%", False, RED_LIGHT, 14),
    ], 7.55, 2.65, 5.6, 3.2, font_size=13, color=LIGHT_GRAY)

    # Result
    add_rect(s, 0.5, 6.25, 12.3, 0.6, fill_color=TEAL)
    add_multiline(s, [
        ("Result: Mann-Whitney U test  p < 0.001  → Reject H₀ — OECD countries have significantly lower PM2.5  "
         "(Δ ≈ 17.4 µg/m³, effect size r = 0.52)", True, HIGHLIGHT, 13),
    ], 0.65, 6.28, 12.0, 0.55)
    add_rect(s, 0.5, 6.9, 12.3, 0.5, fill_color=RGBColor(0x0D, 0x2B, 0x3E))
    add_text(s,
             "Implication: Regulation works — but the 17 µg/m³ gap shows current OECD policy is necessary but not sufficient to reach WHO 15 µg/m³.",
             0.65, 6.93, 12.0, 0.42, font_size=12, color=LIGHT_GRAY)
    return s


def slide_qa(prs):
    s = blank_slide(prs)
    fill_slide_bg(s, DARK_BG)
    add_rect(s, 0, 0, 0.12, 7.5, fill_color=HIGHLIGHT)
    add_rect(s, 0.12, 0, 13.21, 0.08, fill_color=TEAL)

    add_text(s, "Thank You", 0.5, 1.5, 12.3, 1.2,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Questions & Discussion", 0.5, 2.8, 12.3, 0.7,
             font_size=26, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_divider(s, 3.75)

    add_text(s,
             "\"Air pollution is not just an environmental issue — it's a governance failure.\n"
             "The data exists. The solutions exist. What's missing is political will.\"",
             1.0, 4.0, 11.3, 1.5,
             font_size=17, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_multiline(s, [
        ("Oleg Sher  |  claude@sher.biz", False, MID_GRAY, 14),
        ("Data: WHO GHO · IHME GBD · OpenAQ v3 · Open-Meteo ERA5 · World Bank", False, MID_GRAY, 12),
    ], 0.5, 6.6, 12.3, 0.8, align=PP_ALIGN.CENTER)
    return s


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()

    slide_title(prs)
    slide_about_me(prs)
    slide_research_question(prs)
    slide_q1_dirtiest(prs)
    slide_q2_correlation(prs)
    slide_q3_per_capita(prs)
    slide_q4_trends(prs)
    slide_q5_exceedance(prs)
    slide_dashboard(prs)
    slide_conclusions(prs)
    slide_recommendations(prs)
    slide_ab_testing(prs)
    slide_qa(prs)

    out = "data/output/environmental_health_presentation.pptx"
    import os
    os.makedirs("data/output", exist_ok=True)
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
